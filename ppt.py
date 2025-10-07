import streamlit as st
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

def parse_readme_content(content):
    """Parse README.md content and extract sections"""
    sections = []
    current_section = None
    current_content = []
    
    lines = content.split('\n')
    
    for line in lines:
        # Check for main headers (# or ##)
        if line.startswith('# ') or line.startswith('## '):
            # Save previous section
            if current_section:
                sections.append({
                    'title': current_section,
                    'content': '\n'.join(current_content),
                    'level': current_section.count('#')
                })
            
            # Start new section
            current_section = line.strip('#').strip()
            current_content = []
        else:
            if current_section:
                current_content.append(line)
    
    # Add the last section
    if current_section:
        sections.append({
            'title': current_section,
            'content': '\n'.join(current_content),
            'level': current_section.count('#')
        })
    
    return sections

def extract_key_points(content):
    """Extract key points from content"""
    points = []
    
    # Extract bullet points
    bullet_pattern = r'^[\s]*[-*•]\s*(.+)$'
    for line in content.split('\n'):
        match = re.match(bullet_pattern, line)
        if match:
            point = match.group(1).strip()
            # Clean up markdown formatting
            point = re.sub(r'\*\*(.*?)\*\*', r'\1', point)  # Bold
            point = re.sub(r'\*(.*?)\*', r'\1', point)      # Italic
            point = re.sub(r'`(.*?)`', r'\1', point)        # Code
            point = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', point)  # Links
            points.append(point)
    
    # Extract numbered points
    numbered_pattern = r'^[\s]*\d+\.\s*(.+)$'
    for line in content.split('\n'):
        match = re.match(numbered_pattern, line)
        if match:
            point = match.group(1).strip()
            # Clean up markdown formatting
            point = re.sub(r'\*\*(.*?)\*\*', r'\1', point)
            point = re.sub(r'\*(.*?)\*', r'\1', point)
            point = re.sub(r'`(.*?)`', r'\1', point)
            point = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', point)
            points.append(point)
    
    return points[:8]  # Limit to 8 points per slide

def extract_code_blocks(content):
    """Extract code blocks from content"""
    code_blocks = []
    lines = content.split('\n')
    in_code_block = False
    current_code = []
    current_language = ""
    
    for line in lines:
        if line.strip().startswith('```'):
            if in_code_block:
                # End of code block
                code_blocks.append({
                    'language': current_language,
                    'code': '\n'.join(current_code)
                })
                current_code = []
                in_code_block = False
            else:
                # Start of code block
                current_language = line.strip('`').strip() or 'text'
                in_code_block = True
        elif in_code_block:
            current_code.append(line)
    
    return code_blocks

def create_title_slide(prs, title, subtitle=""):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # Set subtitle
    if subtitle and slide.placeholders[1]:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)

def create_content_slide(prs, title, content, slide_type="bullet"):
    """Create content slide with title and bullet points or content"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # Add content
    if slide.placeholders[1]:
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        if slide_type == "bullet" and isinstance(content, list):
            # Add bullet points
            for i, point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(64, 64, 64)
        else:
            # Add regular content
            p = text_frame.paragraphs[0]
            p.text = str(content)
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(64, 64, 64)

def create_table_slide(prs, title, data):
    """Create slide with table"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # Extract table data from markdown-like format
    lines = data.split('\n')
    table_data = []
    for line in lines:
        if '|' in line and not line.strip().startswith('|---'):
            row = [cell.strip() for cell in line.split('|')[1:-1]]
            if row:  # Skip empty rows
                table_data.append(row)
    
    if table_data and len(table_data) > 1:
        rows = len(table_data)
        cols = len(table_data[0])
        
        # Add table
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(5))
        table = table_shape.table
        
        # Fill table with data
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                if j < cols:
                    cell = table.cell(i, j)
                    cell.text = cell_data
                    cell.text_frame.paragraphs[0].font.size = Pt(14)
                    
                    # Header row styling
                    if i == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(31, 73, 125)

def create_code_slide(prs, title, code_block):
    """Create slide with code block"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # Add code block
    code_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    code_frame = code_shape.text_frame
    code_frame.text = code_block['code']
    
    # Style code text
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = 'Consolas'
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    # Add background color to code block
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = RGBColor(248, 248, 248)
    code_shape.line.color.rgb = RGBColor(200, 200, 200)

def generate_ppt_from_readme(readme_content, output_filename):
    """Generate PowerPoint presentation from README content"""
    # Create presentation
    prs = Presentation()
    
    # Parse README content
    sections = parse_readme_content(readme_content)
    
    if not sections:
        st.error("No sections found in README.md")
        return None
    
    # Create title slide
    project_title = "Helix Tag Manager"
    subtitle = "A powerful Django application for migrating Adobe Helix V1 components to V2"
    
    # Find project title from first section
    if sections and sections[0]['title']:
        project_title = sections[0]['title'].replace('🚀', '').strip()
    
    create_title_slide(prs, project_title, subtitle)
    
    # Process each section
    for section in sections:
        title = section['title']
        content = section['content']
        
        # Skip very short sections or navigation sections
        if len(content.strip()) < 50:
            continue
            
        # Extract key points
        key_points = extract_key_points(content)
        
        # Extract code blocks
        code_blocks = extract_code_blocks(content)
        
        # Check if section contains table data
        has_table = '|' in content and '---' in content
        
        # Create appropriate slide type
        if title.lower() in ['api endpoints', 'core components', 'project structure']:
            if has_table:
                create_table_slide(prs, title, content)
            else:
                create_content_slide(prs, title, key_points or ["Content details in documentation"], "bullet")
        elif code_blocks:
            # Create separate slides for code blocks
            for i, code_block in enumerate(code_blocks[:2]):  # Limit to 2 code blocks per section
                code_title = f"{title} - {code_block['language'].title()} Code"
                create_code_slide(prs, code_title, code_block)
            
            # Also create a bullet point slide if there are key points
            if key_points:
                create_content_slide(prs, title, key_points, "bullet")
        elif key_points:
            create_content_slide(prs, title, key_points, "bullet")
        else:
            # Create a simple content slide with cleaned text
            clean_content = re.sub(r'```.*?```', '', content, flags=re.DOTALL)
            clean_content = re.sub(r'\[.*?\]\(.*?\)', '', clean_content)
            clean_content = re.sub(r'[#*`]', '', clean_content)
            lines = [line.strip() for line in clean_content.split('\n') if line.strip()]
            preview_content = ' '.join(lines[:3]) + "..." if len(lines) > 3 else ' '.join(lines)
            
            create_content_slide(prs, title, preview_content, "text")
    
    # Save presentation
    try:
        prs.save(output_filename)
        return output_filename
    except Exception as e:
        st.error(f"Error saving presentation: {str(e)}")
        return None

def main():
    st.set_page_config(
        page_title="README to PPT Converter",
        page_icon="📊",
        layout="wide"
    )
    
    st.title("📊 README to PowerPoint Converter")
    st.markdown("Convert your README.md file into a professional PowerPoint presentation")
    
    # Sidebar configuration
    with st.sidebar:
        st.header("⚙️ Configuration")
        
        # File selection
        readme_path = st.text_input(
            "README.md Path",
            value="/Users/abhi/Desktop/console_script_batch_scrapper/README.md",
            help="Enter the full path to your README.md file"
        )
        
        # Output filename
        output_name = st.text_input(
            "Output Filename",
            value=f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
            help="Name for the generated PowerPoint file"
        )
        
        # Generate button
        generate_button = st.button("🚀 Generate Presentation", type="primary")
    
    # Main content area
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.subheader("📝 README Preview")
        
        # Check if file exists and read it
        if os.path.exists(readme_path):
            try:
                with open(readme_path, 'r', encoding='utf-8') as file:
                    readme_content = file.read()
                
                # Show preview
                st.text_area(
                    "README Content Preview",
                    value=readme_content[:2000] + "..." if len(readme_content) > 2000 else readme_content,
                    height=400,
                    disabled=True
                )
                
                # Show section count
                sections = parse_readme_content(readme_content)
                st.info(f"📊 Found {len(sections)} sections in README.md")
                
                # Show sections list
                if sections:
                    st.subheader("📋 Sections to be included:")
                    for i, section in enumerate(sections[:10], 1):  # Show first 10 sections
                        st.write(f"{i}. {section['title']}")
                    
                    if len(sections) > 10:
                        st.write(f"... and {len(sections) - 10} more sections")
                        
            except Exception as e:
                st.error(f"Error reading README.md: {str(e)}")
                readme_content = None
        else:
            st.error(f"README.md file not found at: {readme_path}")
            readme_content = None
    
    with col2:
        st.subheader("🎯 Generation Status")
        
        if generate_button:
            if readme_content:
                with st.spinner("🔄 Generating PowerPoint presentation..."):
                    # Generate PPT
                    output_path = os.path.join(
                        os.path.dirname(readme_path),
                        output_name
                    )
                    
                    result = generate_ppt_from_readme(readme_content, output_path)
                    
                    if result:
                        st.success("✅ Presentation generated successfully!")
                        st.balloons()
                        
                        # Show file info
                        file_size = os.path.getsize(result) / (1024 * 1024)  # MB
                        st.info(f"📁 File saved as: {result}")
                        st.info(f"📊 File size: {file_size:.2f} MB")
                        
                        # Download button
                        with open(result, 'rb') as file:
                            st.download_button(
                                label="📥 Download Presentation",
                                data=file.read(),
                                file_name=output_name,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                    else:
                        st.error("❌ Failed to generate presentation")
            else:
                st.error("❌ Please provide a valid README.md file")
        
        # Instructions
        st.subheader("📖 How to use:")
        st.markdown("""
        1. **Enter README path**: Provide the full path to your README.md file
        2. **Set output name**: Choose a name for your presentation file
        3. **Click Generate**: The tool will create a PowerPoint presentation
        4. **Download**: Use the download button to get your presentation
        
        **Features:**
        - 🎨 Professional slide design
        - 📊 Automatic table conversion
        - 💻 Code block formatting
        - 📋 Bullet point extraction
        - 🎯 Section-based organization
        """)
        
        # Requirements info
        with st.expander("📦 Required Dependencies"):
            st.code("""
            pip install streamlit python-pptx
            """)

if __name__ == "__main__":
    main()